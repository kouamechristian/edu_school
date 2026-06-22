# -*- coding: utf-8 -*-
"""Génère la présentation PowerPoint d'EDU-SCHOOL, module par module."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Charte graphique (Navy + Émeraude) ──
NAVY      = RGBColor(0x0F, 0x17, 0x2A)
NAVY_2    = RGBColor(0x1E, 0x29, 0x3B)
EMERALD   = RGBColor(0x10, 0xB9, 0x81)
EMERALD_D = RGBColor(0x05, 0x96, 0x69)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xE2, 0xE8, 0xF0)
GREY      = RGBColor(0x64, 0x74, 0x8B)
SLATE     = RGBColor(0x33, 0x41, 0x55)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, color, line=False):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = color
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size, color, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri", italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font
    return tb


def bullets(slide, x, y, w, h, items, size=15, color=SLATE, gap=6, bullet_color=EMERALD):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.level = lvl
        mark = "▸  " if lvl == 0 else "•  "
        r1 = p.add_run(); r1.text = mark
        r1.font.size = Pt(size); r1.font.bold = True
        r1.font.color.rgb = bullet_color if lvl == 0 else GREY
        r1.font.name = "Calibri"
        r2 = p.add_run(); r2.text = txt
        r2.font.size = Pt(size if lvl == 0 else size - 1)
        r2.font.color.rgb = color
        r2.font.bold = (lvl == 0)
        r2.font.name = "Calibri"
    return tb


def content_header(slide, kicker, title):
    """Bandeau supérieur des slides de contenu."""
    add_rect(slide, 0, 0, SW, Inches(1.15), NAVY)
    add_rect(slide, 0, Inches(1.15), SW, Pt(4), EMERALD)
    add_rect(slide, Inches(0.0), Inches(0.18), Inches(0.16), Inches(0.8), EMERALD)
    add_text(slide, Inches(0.45), Inches(0.16), Inches(11), Inches(0.35),
             kicker, 12, EMERALD, bold=True)
    add_text(slide, Inches(0.45), Inches(0.45), Inches(12.4), Inches(0.62),
             title, 26, WHITE, bold=True)


def footer(slide, n):
    add_text(slide, Inches(0.45), Inches(7.05), Inches(6), Inches(0.3),
             "EDU-SCHOOL — Plateforme de gestion scolaire", 9, GREY)
    add_text(slide, Inches(11.8), Inches(7.05), Inches(1.1), Inches(0.3),
             str(n), 9, GREY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
# 1. SLIDE DE TITRE
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, 0, Inches(0.35), SH, EMERALD)
# pastille
add_text(s, Inches(1.0), Inches(1.6), Inches(11), Inches(0.5),
         "GESTION D'ÉTABLISSEMENTS SCOLAIRES", 16, EMERALD, bold=True)
add_text(s, Inches(0.95), Inches(2.15), Inches(11.5), Inches(1.6),
         "EDU-SCHOOL", 76, WHITE, bold=True)
add_text(s, Inches(1.0), Inches(3.7), Inches(11), Inches(0.8),
         "Plateforme web multi-établissements & multi-années",
         24, LIGHT)
add_rect(s, Inches(1.0), Inches(4.6), Inches(3.2), Pt(3), EMERALD)
add_text(s, Inches(1.0), Inches(4.85), Inches(11.5), Inches(0.5),
         "Symfony 6.4  ·  PHP 8.1+  ·  Doctrine ORM 3  ·  MySQL  ·  Twig / Bootstrap  ·  IA Claude",
         14, GREY)
add_text(s, Inches(1.0), Inches(6.6), Inches(11), Inches(0.4),
         "Présentation du projet — module par module", 14, EMERALD, bold=True)

# ════════════════════════════════════════════════════════════
# 2. SOMMAIRE
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
content_header(s, "VUE D'ENSEMBLE", "Sommaire")
col1 = [
    ("1.  Présentation générale", 0),
    ("2.  Architecture & stack technique", 0),
    ("3.  Sécurité & rôles", 0),
    ("4.  Contexte établissement / année", 0),
    ("5.  Administration & Configuration", 0),
    ("6.  Gestion des élèves", 0),
    ("7.  Académique", 0),
    ("8.  Notes & Évaluations", 0),
]
col2 = [
    ("9.   Absences & Assiduité", 0),
    ("10.  Ressources Humaines", 0),
    ("11.  Finances & Caisse", 0),
    ("12.  Recouvrement", 0),
    ("13.  Paiement en ligne (GeniusPay)", 0),
    ("14.  Espaces Parent & Fondateur", 0),
    ("15.  Communication & Rapports", 0),
    ("16.  Intelligence Artificielle", 0),
]
bullets(s, Inches(0.7), Inches(1.6), Inches(6), Inches(5.2), col1, size=17, gap=12)
bullets(s, Inches(6.9), Inches(1.6), Inches(6), Inches(5.2), col2, size=17, gap=12)
footer(s, 2)

# ════════════════════════════════════════════════════════════
# 3. PRÉSENTATION GÉNÉRALE
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
content_header(s, "01 · INTRODUCTION", "Présentation générale")
add_text(s, Inches(0.7), Inches(1.45), Inches(12), Inches(0.9),
         "EDU-SCHOOL couvre tout le cycle de vie de l'élève — de la préinscription à la "
         "scolarité, aux notes/bulletins et aux finances — ainsi que l'administration de "
         "l'établissement (RH, caisse, recouvrement, communication).", 15, SLATE)
# KPI cards
kpis = [("44", "Contrôleurs"), ("41", "Entités"), ("29", "Services"),
        ("34", "Formulaires"), ("209", "Templates"), ("13", "Rôles")]
cx = Inches(0.7); cy = Inches(2.7)
cw = Inches(1.9); ch = Inches(1.5); gap = Inches(0.12)
for i, (val, lbl) in enumerate(kpis):
    x = Emu(int(cx) + i * (int(cw) + int(gap)))
    card = add_rect(s, x, cy, cw, ch, NAVY)
    add_text(s, x, Inches(2.95), cw, Inches(0.7), val, 34, EMERALD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(3.75), cw, Inches(0.4), lbl, 12, LIGHT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(4.6), Inches(12), Inches(0.4),
         "Principes fondateurs", 18, NAVY, bold=True)
bullets(s, Inches(0.7), Inches(5.05), Inches(12), Inches(1.9), [
    ("Multi-établissement : cloisonnement automatique des données par établissement (filtre Doctrine).", 0),
    ("Multi-année : la scolarité est portée par une inscription annuelle, pas par l'élève.", 0),
    ("Student = référentiel d'identité pur ; toute la scolarité vit dans Registration.", 0),
    ("Sécurité par rôles hiérarchiques (du parent au super-administrateur).", 0),
], size=14, gap=7)
footer(s, 3)

# ════════════════════════════════════════════════════════════
# 4. ARCHITECTURE & STACK
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
content_header(s, "02 · TECHNIQUE", "Architecture & stack technique")
# deux colonnes en cartes
add_rect(s, Inches(0.7), Inches(1.5), Inches(5.9), Inches(2.5), NAVY_2)
add_text(s, Inches(0.95), Inches(1.65), Inches(5.4), Inches(0.4), "Socle technique", 16, EMERALD, bold=True)
bullets(s, Inches(0.95), Inches(2.2), Inches(5.5), Inches(1.8), [
    ("Symfony 6.4 (LTS) · PHP 8.1+", 0),
    ("Doctrine ORM 3 + migrations · MySQL", 0),
    ("Twig + Bootstrap 5 + FontAwesome", 0),
    ("VichUploader · KnpPaginator (50/page)", 0),
], size=13, gap=5, color=LIGHT, bullet_color=EMERALD)

add_rect(s, Inches(6.75), Inches(1.5), Inches(5.9), Inches(2.5), NAVY_2)
add_text(s, Inches(7.0), Inches(1.65), Inches(5.4), Inches(0.4), "Briques métier", 16, EMERALD, bold=True)
bullets(s, Inches(7.0), Inches(2.2), Inches(5.5), Inches(1.8), [
    ("Dompdf — reçus, bulletins, rapports PDF", 0),
    ("PhpSpreadsheet — imports/exports Excel", 0),
    ("endroid/qr-code — reçus & badges", 0),
    ("CKEditor + elFinder — documents", 0),
    ("API Anthropic Claude · GeniusPay", 0),
], size=13, gap=5, color=LIGHT, bullet_color=EMERALD)

add_text(s, Inches(0.7), Inches(4.3), Inches(12), Inches(0.4),
         "Organisation du code (src/)", 18, NAVY, bold=True)
bullets(s, Inches(0.7), Inches(4.8), Inches(12), Inches(2), [
    ("Controller/ — 44 contrôleurs (+ Concern, Portal, Webhook)", 0),
    ("Entity/ — 41 entités Doctrine", 0),
    ("Service/ — 29 services métier (+ AI/, Payment/)", 0),
    ("Form/ · Repository/ · Security/ · EventSubscriber/ · Doctrine/Filter/", 0),
], size=14, gap=6)
footer(s, 4)

# ════════════════════════════════════════════════════════════
# 5. SÉCURITÉ & RÔLES
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
content_header(s, "03 · SÉCURITÉ", "Sécurité & rôles hiérarchiques")
add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
         "Connexion par identifiant OU e-mail · anti-bruteforce (3 essais/2 min) · « se souvenir de moi ».",
         13, GREY, italic=True)
# pyramide simplifiée
levels = [
    ("ROLE_SUPER_ADMIN", "Sommet de la hiérarchie", Inches(4.0)),
    ("ROLE_FONDATEUR", "Supervision de ses établissements", Inches(5.0)),
    ("ROLE_ADMIN", "Administration complète d'un établissement", Inches(6.0)),
    ("ROLE_DIRECTEUR / CAISSE / INSCRIPTION / RH", "Rôles métier", Inches(8.5)),
]
y = Inches(2.05)
for name, desc, w in levels:
    x = Emu(int((SW - w) / 2))
    add_rect(s, x, y, w, Inches(0.62), NAVY if "SUPER" not in name else EMERALD_D)
    add_text(s, x, y + Inches(0.06), w, Inches(0.5), name, 13, WHITE, bold=True, align=PP_ALIGN.CENTER)
    y = Emu(int(y) + int(Inches(0.72)))
add_text(s, Inches(0.7), Inches(5.25), Inches(12), Inches(0.4),
         "Rôles opérationnels", 15, NAVY, bold=True)
bullets(s, Inches(0.7), Inches(5.65), Inches(6), Inches(1.5), [
    ("ENSEIGNANT — ses cours, élèves, évaluations", 0),
    ("EDUCATEUR — saisie/suivi des absences", 0),
    ("INSCRIPTION — gestion des élèves", 0),
], size=12.5, gap=4)
bullets(s, Inches(6.9), Inches(5.65), Inches(6), Inches(1.5), [
    ("CAISSE — finances (hérite INSCRIPTION + RECOUVREMENT)", 0),
    ("RECOUVREMENT — soldes & relances", 0),
    ("PARENT — espace parent uniquement", 0),
], size=12.5, gap=4)
footer(s, 5)

# ════════════════════════════════════════════════════════════
# 6. CONTEXTE ÉTABLISSEMENT / ANNÉE
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
content_header(s, "04 · TRANSVERSE", "Contexte établissement & année")
bullets(s, Inches(0.7), Inches(1.6), Inches(12), Inches(4.5), [
    ("Établissement courant + année scolaire courante mémorisés en session.", 0),
    ("SchoolContextService — getCurrentSchool() / setCurrentSchool() / getCurrentSchoolYear().", 1),
    ("Un utilisateur peut être rattaché à plusieurs établissements (ManyToMany) et basculer.", 0),
    ("Le dernier établissement utilisé est persisté sur le compte (lastSchool).", 1),
    ("SchoolFilter (filtre Doctrine) applique automatiquement la restriction par établissement.", 0),
    ("Garantit le cloisonnement des données entre établissements.", 1),
    ("SchoolContextSubscriber active le contexte à chaque requête.", 0),
], size=15, gap=10)
footer(s, 6)

# ════════════════════════════════════════════════════════════
# Helper : slide de module standard
# ════════════════════════════════════════════════════════════
def module_slide(n, kicker, title, role, entities, controllers, rules):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, WHITE)
    content_header(s, kicker, title)
    # badge rôle
    badge = add_rect(s, Inches(0.7), Inches(1.4), Inches(5.5), Inches(0.5), EMERALD)
    add_text(s, Inches(0.85), Inches(1.47), Inches(5.3), Inches(0.4),
             "Accès : " + role, 12, WHITE, bold=True)
    # colonne gauche : entités
    add_text(s, Inches(0.7), Inches(2.15), Inches(6), Inches(0.4), "Entités & données", 15, NAVY, bold=True)
    bullets(s, Inches(0.7), Inches(2.6), Inches(6), Inches(3.0), entities, size=12.5, gap=5)
    # colonne droite : contrôleurs/services
    add_text(s, Inches(6.95), Inches(2.15), Inches(6), Inches(0.4), "Contrôleurs & services", 15, NAVY, bold=True)
    bullets(s, Inches(6.95), Inches(2.6), Inches(6), Inches(3.0), controllers, size=12.5, gap=5)
    # bandeau règles métier
    add_rect(s, Inches(0.7), Inches(5.75), Inches(11.95), Inches(1.1), NAVY_2)
    add_text(s, Inches(0.9), Inches(5.83), Inches(11.5), Inches(0.35), "Règles métier clés", 13, EMERALD, bold=True)
    bullets(s, Inches(0.9), Inches(6.18), Inches(11.6), Inches(0.7), rules, size=11.5, gap=2, color=LIGHT, bullet_color=EMERALD)
    footer(s, n)
    return s


# ── 7. Administration & Configuration ──
module_slide(7, "05 · STRUCTURE", "Administration & Configuration", "ROLE_ADMIN",
    [("School — fiche établissement (type, logo, cachet, tutelle)", 0),
     ("SchoolGroup — réseau d'établissements", 0),
     ("SchoolYear → Period — années & trimestres/semestres", 0),
     ("Cycle → Level → Subject — structure pédagogique", 0),
     ("Faculty, Round, Room, TimeSlot — filières, séries, salles, créneaux", 0),
     ("SubjectType, DocumentType — référentiels", 0)],
    [("SchoolController · SchoolGroupController", 0),
     ("SchoolYearController · UserController", 0),
     ("CycleController · LevelController · FacultyController", 0),
     ("SubjectController · RoomController · TimeSlotController", 0),
     ("RoundController · DocumentTypeController", 0)],
    [("Hiérarchie pédagogique : Cycle → Level → Classroom ; Subject rattaché à un Level.", 0)])

# ── 8. Gestion des élèves ──
module_slide(8, "06 · ÉLÈVES", "Gestion des élèves", "ROLE_INSCRIPTION",
    [("Student — référentiel d'identité (état civil, contacts, matricules)", 0),
     ("PreRegistration — demande d'inscription (pending → validé → inscrit)", 0),
     ("PreRegistrationDocument — pièces jointes", 0),
     ("Registration — inscription annuelle (École + Année + Classe)", 0),
     ("StudentDropout / StudentTransfer — abandons & transferts", 0)],
    [("PreRegistrationController · RegistrationController", 0),
     ("StudentController · StudentDropoutController", 0),
     ("EnrollmentService — inscription depuis préinscription", 0),
     ("RegistrationManager — source de vérité scolarité", 0),
     ("MatriculeGenerator · PreRegistrationFactory", 0)],
    [("Workflow : Préinscription → Validation → Inscription. Une seule inscription par élève et par année ; ancien élève réutilisé sans duplication.", 0)])

# ── 9. Académique ──
module_slide(9, "07 · ACADÉMIQUE", "Académique — classes & emploi du temps", "ROLE_DIRECTEUR",
    [("Classroom — classe d'une année (niveau, filière, capacité, prof principal, salle)", 0),
     ("Course — matière + enseignant + créneau + salle", 0),
     ("Period — trimestres / semestres", 0)],
    [("ClassroomController · CourseController · PeriodController", 0),
     ("Espace Enseignant :", 0),
     ("Mes Cours (emploi du temps)", 1),
     ("Mes élèves (classes enseignées)", 1),
     ("Évaluations (saisie des notes)", 1)],
    [("Le niveau d'une inscription se déduit de la classe ; l'emploi du temps relie Classroom, Subject, TimeSlot et Room.", 0)])

# ── 10. Notes & Évaluations ──
module_slide(10, "08 · ÉVALUATION", "Notes & Évaluations", "ROLE_ENSEIGNANT / DIRECTEUR",
    [("Evaluation — devoir/composition (barème, coefficient, période)", 0),
     ("Grade — note d'un élève (valeur, statut, appréciation)", 0),
     ("GeneratedBulletin — trace des bulletins générés", 0)],
    [("EvaluationController · BulletinController", 0),
     ("AcademicReportController — stats, majors, répartitions", 0),
     ("GradeCalculationService :", 0),
     ("moyennes pondérées par coefficient", 1),
     ("rang de l'élève, bulletin officiel + mention", 1)],
    [("Bulletins produits en PDF (Dompdf) ; appréciations pouvant être générées par IA (BulletinAIService).", 0)])

# ── 11. Absences ──
module_slide(11, "09 · ASSIDUITÉ", "Absences & Assiduité", "ROLE_EDUCATEUR",
    [("Absence — date, plage horaire, motif, justification", 0),
     ("AbsenceType — catégories d'absence", 0),
     ("Workflow justification : pending → validé (+ justificatif)", 0)],
    [("AbsenceController — saisie & justification", 0),
     ("AttendanceService :", 0),
     ("stats par élève / classe / établissement", 1),
     ("taux de présence, rapport détaillé", 1),
     ("détection assiduité critique (seuil 75 %)", 1)],
    [("Le taux d'assiduité alimente les bulletins, les rapports et les alertes de l'espace parent.", 0)])

# ── 12. RH ──
module_slide(12, "10 · RH", "Ressources Humaines", "ROLE_RH",
    [("Employee — dossier (type, poste, département, salaire, dates)", 0),
     ("Contract — contrat (type, période d'essai, salaire, heures/sem, statut)", 0),
     ("Teacher — profil enseignant lié à un employé", 0)],
    [("EmployeeController · ContractController · TeacherController", 0),
     ("UserEmployeeService — synchronise compte ↔ fiche employé", 0)],
    [("Un employé peut être lié à un compte utilisateur et à un profil enseignant ; il possède 0..n contrats.", 0)])

# ── 13. Finances ──
module_slide(13, "11 · FINANCES", "Finances — Caisse & Frais", "ROLE_CAISSE",
    [("Fee — rubrique (catégorie, type, fréquence) + FeeSchedule (échéances)", 0),
     ("StudentFee — frais affecté à l'élève (dû / payé / statut)", 0),
     ("Payment — paiement (méthode, statut, reçu PDF)", 0),
     ("CashRegister — caisse · CashDeposit — versement", 0),
     ("Depense — dépense imputée à une caisse", 0)],
    [("FeeController · PaymentController", 0),
     ("CashRegisterController · DepenseController", 0),
     ("FeeAssignmentService — affectation/sync des frais", 0),
     ("PaymentReceiptService — reçu PDF détaillé", 0)],
    [("Solde caisse = Paiements − Versements approuvés − Dépenses · Reçu en toutes lettres (francs CFA).", 0)])

# ── 14. Recouvrement ──
module_slide(14, "12 · RECOUVREMENT", "Recouvrement & relances", "ROLE_RECOUVREMENT",
    [("Lignes de recouvrement par élève : dû, payé, solde, statut", 0),
     ("Montant échu impayé, nb d'échéances échues, jours de retard", 0),
     ("Prochaine échéance (date + montant)", 0)],
    [("RecouvrementController — tableau de bord des soldes", 0),
     ("RecouvrementService :", 0),
     ("filtrage par catégorie de frais et par année", 1),
     ("statut de relance basé sur les échéanciers", 1)],
    [("Seules les échéances dépassées et non couvertes par les paiements justifient une relance.", 0)])

# ── 15. Paiement en ligne (GeniusPay) ──
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
content_header(s, "13 · PAIEMENT EN LIGNE", "Paiement en ligne — GeniusPay (Mobile Money)")
add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
         "Les parents règlent les frais en ligne via la passerelle GeniusPay (Mobile Money).", 14, SLATE)
# flux horizontal
steps = ["Parent\n(espace parent)", "PaymentInitiator\ninitiate()", "GeniusPay\ncheckoutUrl",
         "Webhook signé\n/webhook", "WebhookProcessor\nMAJ Payment + StudentFee"]
x = Inches(0.7); y = Inches(2.2); bw = Inches(2.25); bh = Inches(1.0); g = Inches(0.16)
for i, st in enumerate(steps):
    bx = Emu(int(x) + i * (int(bw) + int(g)))
    add_rect(s, bx, y, bw, bh, NAVY if i % 2 == 0 else EMERALD_D)
    add_text(s, bx, y + Inches(0.12), bw, Inches(0.8), st, 11.5, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.7), Inches(3.7), Inches(12), Inches(0.4), "Composants", 16, NAVY, bold=True)
bullets(s, Inches(0.7), Inches(4.15), Inches(6), Inches(2.5), [
    ("GeniusPayClient — createPayment() / getPaymentStatus()", 0),
    ("GeniusPaySignatureVerifier — vérifie les webhooks", 0),
    ("PaymentStatusSynchronizer — synchronise les statuts", 0),
    ("OnlineCashRegisterProvider — caisse virtuelle", 0),
], size=13, gap=6)
bullets(s, Inches(6.95), Inches(4.15), Inches(6), Inches(2.5), [
    ("MobileMoneyConfig — identifiants par établissement", 0),
    ("PaymentWebhookEvent — journal des callbacks", 0),
    ("Payment : provider, providerTransactionId, checkoutUrl, idempotencyKey", 0),
    ("Webhook public sécurisé par signature", 0),
], size=13, gap=6)
footer(s, 15)

# ── 16. Espaces Parent & Fondateur ──
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
content_header(s, "14 · ESPACES DÉDIÉS", "Espaces Parent & Fondateur")
add_rect(s, Inches(0.7), Inches(1.5), Inches(5.9), Inches(4.9), NAVY_2)
add_text(s, Inches(0.95), Inches(1.65), Inches(5.4), Inches(0.45), "Espace Parent  ·  ROLE_PARENT", 15, EMERALD, bold=True)
bullets(s, Inches(0.95), Inches(2.25), Inches(5.5), Inches(4.0), [
    ("ParentPortalService — tableau de bord par enfant/année", 0),
    ("getChildren() — enfants rattachés", 1),
    ("getAcademicReport() — notes & moyennes", 1),
    ("getAttendanceReport() — assiduité", 1),
    ("getFinancialReport() — situation financière", 1),
    ("Réinscription en ligne d'un ancien élève (→ préinscription pending)", 0),
    ("Paiement en ligne des frais (GeniusPay)", 0),
], size=13, gap=6, color=LIGHT, bullet_color=EMERALD)
add_rect(s, Inches(6.75), Inches(1.5), Inches(5.9), Inches(4.9), NAVY_2)
add_text(s, Inches(7.0), Inches(1.65), Inches(5.4), Inches(0.45), "Espace Fondateur  ·  ROLE_FONDATEUR", 15, EMERALD, bold=True)
bullets(s, Inches(7.0), Inches(2.25), Inches(5.5), Inches(4.0), [
    ("Supervision multi-établissements", 0),
    ("Tableau de bord global", 1),
    ("Validation des caisses (clôtures)", 1),
    ("Gestion des autorisations", 1),
    ("Validation des versements (CashDeposit)", 1),
    ("Périmètre limité à ses propres établissements", 0),
], size=13, gap=6, color=LIGHT, bullet_color=EMERALD)
footer(s, 16)

# ── 17. Communication & Rapports ──
module_slide(17, "15 · PILOTAGE", "Communication & Rapports", "Selon le rôle",
    [("Notification — alertes internes (titre, message, lien, icône)", 0),
     ("Badge non-lus dans la barre supérieure, marquage lu", 0),
     ("E-mails via Symfony Mailer", 0)],
    [("NotificationController · NotificationService", 0),
     ("notify() — notification individuelle", 1),
     ("notifyRole() — diffusion à tous les rôles ciblés", 1),
     ("ReportController · AcademicReportController", 0),
     ("Exports PDF (Dompdf) et Excel (PhpSpreadsheet)", 0)],
    [("Les rapports peuvent être synthétisés par IA (ReportAIService) ; documents gérés via elFinder/CKEditor.", 0)])

# ── 18. Intelligence Artificielle ──
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
content_header(s, "16 · IA", "Intelligence Artificielle (API Claude)")
add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
         "Intégration de l'API Anthropic Claude — activable via AI_ENABLED, modèle configurable (AI_MODEL).",
         14, SLATE)
add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.4), "Service de base — AIService", 16, NAVY, bold=True)
bullets(s, Inches(0.7), Inches(2.45), Inches(12), Inches(1.0), [
    ("Appelle l'API Claude · ask() avec cache (AI_CACHE_TTL) / askWithoutCache() · plafond AI_MAX_TOKENS.", 0),
], size=13, gap=4)
add_text(s, Inches(0.7), Inches(3.35), Inches(12), Inches(0.4), "Services spécialisés", 16, NAVY, bold=True)
cards = [
    ("BulletinAIService", "Appréciations de bulletin\nà partir des notes"),
    ("AttendanceAIService", "Analyse des absences\n+ recommandations"),
    ("ReportAIService", "Synthèse des rapports\nd'établissement"),
    ("ChatbotAIService", "Assistant conversationnel\ncontextualisé par profil"),
]
x = Inches(0.7); y = Inches(3.85); cw = Inches(2.95); ch = Inches(2.0); g = Inches(0.13)
for i, (t, d) in enumerate(cards):
    bx = Emu(int(x) + i * (int(cw) + int(g)))
    add_rect(s, bx, y, cw, ch, NAVY)
    add_rect(s, bx, y, cw, Inches(0.12), EMERALD)
    add_text(s, bx, y + Inches(0.35), cw, Inches(0.6), t, 14, EMERALD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, bx, y + Inches(1.0), cw, Inches(0.9), d, 12, LIGHT, align=PP_ALIGN.CENTER)
footer(s, 18)

# ── 19. SLIDE DE CLÔTURE ──
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, 0, Inches(0.35), SH, EMERALD)
add_text(s, Inches(1.0), Inches(2.5), Inches(11), Inches(1.2), "Merci", 60, WHITE, bold=True)
add_rect(s, Inches(1.05), Inches(3.8), Inches(3.0), Pt(3), EMERALD)
add_text(s, Inches(1.0), Inches(4.1), Inches(11), Inches(0.6),
         "EDU-SCHOOL — une plateforme scolaire complète, modulaire et multi-établissements.",
         18, LIGHT)
add_text(s, Inches(1.0), Inches(5.0), Inches(11), Inches(0.5),
         "Préinscription · Scolarité · Notes & Bulletins · Finances · RH · IA",
         14, EMERALD, bold=True)

prs.save(r"C:\xampp\htdocs\edu_school\EDU-SCHOOL_Presentation.pptx")
print("OK -", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
